from pptx import Presentation
import sys

def extract_text(pptx_path):
    try:
        prs = Presentation(pptx_path)
        for i, slide in enumerate(prs.slides):
            sys.stdout.write(f"\n--- SLIDE {i+1} ---\n")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip().replace("\n", " ")
                    sys.stdout.write(f"[{shape.name}] {text}\n")
    except Exception as e:
        sys.stdout.write(f"Error: {e}\n")

if __name__ == "__main__":
    extract_text(r"e:\SwarajyaRatna\SwarajyaRatna Startup.pptx")
